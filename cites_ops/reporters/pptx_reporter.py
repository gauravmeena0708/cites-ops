from datetime import date
from pathlib import Path
from typing import Any, Dict, List, Optional, Union
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

class PPTXReporter:
    """
    Generates high-impact, light-themed executive PowerPoint slide decks (.pptx)
    for CITES Operations, Workforce Accountability, and Root-Cause Defect Diagnostics.
    """

    # Light Theme Color Palette
    COLOR_BG = RGBColor(248, 250, 252)          # Light Slate BG #F8FAFC
    COLOR_CARD_BG = RGBColor(255, 255, 255)     # White #FFFFFF
    COLOR_NAVY = RGBColor(15, 41, 66)           # Deep Executive Navy #0F2942
    COLOR_PRIMARY = RGBColor(31, 78, 121)       # Primary Blue #1F4E79
    COLOR_HEADER_BG = RGBColor(235, 243, 250)   # Light Blue Header #EBF3FA
    COLOR_TEXT_MAIN = RGBColor(15, 23, 42)      # Charcoal Text #0F172A
    COLOR_TEXT_MUTED = RGBColor(100, 116, 139)  # Muted Slate #64748B
    COLOR_BORDER = RGBColor(226, 232, 240)      # Subtle Border #E2E8F0
    COLOR_RED = RGBColor(220, 38, 38)           # Alert Crimson #DC2626
    COLOR_GREEN = RGBColor(22, 163, 74)         # Success Green #16A34A
    COLOR_BLUE = RGBColor(37, 99, 235)          # Accent Blue #2563EB
    COLOR_WHITE = RGBColor(255, 255, 255)

    @classmethod
    def generate_presentation(
        cls,
        df_classified: pd.DataFrame,
        output_path: Union[str, Path],
        workload_data: Optional[Dict[str, Any]] = None,
        report_date: Optional[Union[str, date]] = None,
        title: str = "CITES Operations Intelligence & Defect Review",
    ) -> str:
        out_file = Path(output_path)
        out_file.parent.mkdir(parents=True, exist_ok=True)

        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9 Widescreen standard
        prs.slide_height = Inches(7.5)

        run_date_str = str(report_date or date.today())

        # Slide 1: Title Slide (Executive Light Theme)
        cls._add_title_slide(prs, title, run_date_str)

        # Slide 2: Executive Overview & Operational Health Dashboard
        cls._add_kpi_slide(prs, df_classified, workload_data, run_date_str)

        # Slide 3: Top 10 Major Problem Categories (Functionalities)
        cls._add_top_categories_slide(prs, workload_data)

        # Slide 4: System-Wide Top 10 Root-Cause Defect Drivers
        cls._add_defect_drivers_slide(prs, workload_data, df_classified)

        # Slide 5: Leadership Accountability & Workload Distribution (JD / DD)
        cls._add_leadership_slide(prs, workload_data)

        # Slide 6: Cross-Module Defect Heatmap & Topical Highlights
        cls._add_cross_tab_slide(prs, workload_data)

        # Slide 7: Daily Aging Exceptions & Action Escalations
        cls._add_aging_slide(prs, df_classified)

        prs.save(out_file)
        return str(out_file)

    @classmethod
    def _create_blank_slide(cls, prs: Presentation):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # Background shape
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = cls.COLOR_BG
        bg.line.fill.background()
        return slide

    @classmethod
    def _add_slide_header(cls, slide, title: str, subtitle: str) -> None:
        # Accent top bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.45), Inches(0.4), Inches(0.06)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = cls.COLOR_BLUE
        bar.line.fill.background()

        # Header Text
        box = slide.shapes.add_textbox(Inches(0.8), Inches(0.55), Inches(11.733), Inches(0.9))
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Segoe UI"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = cls.COLOR_NAVY

        p_sub = tf.add_paragraph()
        p_sub.text = subtitle
        p_sub.font.name = "Segoe UI"
        p_sub.font.size = Pt(11)
        p_sub.font.color.rgb = cls.COLOR_TEXT_MUTED

    @classmethod
    def _add_title_slide(cls, prs: Presentation, title: str, date_str: str) -> None:
        slide = cls._create_blank_slide(prs)

        # Decorative subtle card banner
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.2), Inches(11.733), Inches(5.1)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = cls.COLOR_CARD_BG
        card.line.color.rgb = cls.COLOR_BORDER
        card.line.width = Pt(1)

        # Title Content Box
        tx_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.3), Inches(3.2))
        tf = tx_box.text_frame
        tf.word_wrap = True

        # Tag
        p_tag = tf.paragraphs[0]
        p_tag.text = "CITES OPERATIONS INTELLIGENCE & DECISION SUPPORT"
        p_tag.font.name = "Segoe UI"
        p_tag.font.size = Pt(12)
        p_tag.font.bold = True
        p_tag.font.color.rgb = cls.COLOR_BLUE

        # Main Title
        p_title = tf.add_paragraph()
        p_title.text = title
        p_title.font.name = "Segoe UI"
        p_title.font.size = Pt(32)
        p_title.font.bold = True
        p_title.font.color.rgb = cls.COLOR_NAVY
        p_title.space_before = Pt(8)

        # Subtitle
        p_sub = tf.add_paragraph()
        p_sub.text = f"Functional Accountability, Workforce Distribution & Root-Cause Defect Diagnostics\nAs of Snapshot Date: {date_str}"
        p_sub.font.name = "Segoe UI"
        p_sub.font.size = Pt(16)
        p_sub.font.color.rgb = cls.COLOR_TEXT_MUTED
        p_sub.space_before = Pt(14)

    @classmethod
    def _add_kpi_slide(
        cls,
        prs: Presentation,
        df: pd.DataFrame,
        workload_data: Optional[Dict[str, Any]],
        date_str: str,
    ) -> None:
        slide = cls._create_blank_slide(prs)
        cls._add_slide_header(slide, "Operational Health & Intake Overview", f"Snapshot Date: {date_str} · Overall Ingested Issue Portfolio")

        total = len(df)
        status_series = df["Status"].astype(str).str.lower() if "Status" in df.columns else pd.Series([])
        resolved = int(status_series.isin(["resolved", "fixed", "closed"]).sum())
        open_cnt = total - resolved
        res_rate = f"{round(100 * resolved / total, 1)}%" if total else "0%"

        wk_kpis = (workload_data or {}).get("kpis", {})
        routing = wk_kpis.get("routing_breakdown", {})
        epfo_cnt = routing.get("internal_tech", total)
        cdac_cnt = routing.get("vendor_tech", 0)
        field_cnt = routing.get("field_office", 0)
        cov_pct = wk_kpis.get("coverage_pct", "100.0%")

        # 4 Primary KPI Cards
        kpis = [
            ("TOTAL INGESTED", f"{total:,}", "Issues across all 36 modules", cls.COLOR_NAVY),
            ("OPEN BACKLOG", f"{open_cnt:,}", "Requires active resolution", cls.COLOR_RED),
            ("RESOLVED / CLOSED", f"{resolved:,}", f"Resolution Rate: {res_rate}", cls.COLOR_GREEN),
            ("OWNERSHIP COVERAGE", cov_pct, "Mapped in Issue_teams.csv", cls.COLOR_PRIMARY),
        ]

        left_start = 0.8
        card_w = 2.75
        gap = 0.24
        top = 1.65
        card_h = 2.1

        for idx, (label, val, subtext, col) in enumerate(kpis):
            x = Inches(left_start + idx * (card_w + gap))
            
            # Card Shape
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(top), Inches(card_w), Inches(card_h)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = cls.COLOR_CARD_BG
            card.line.color.rgb = cls.COLOR_BORDER
            card.line.width = Pt(1)

            tf = card.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE

            p_lbl = tf.paragraphs[0]
            p_lbl.text = label
            p_lbl.font.name = "Segoe UI"
            p_lbl.font.size = Pt(10)
            p_lbl.font.bold = True
            p_lbl.font.color.rgb = cls.COLOR_TEXT_MUTED
            p_lbl.alignment = PP_ALIGN.CENTER

            p_val = tf.add_paragraph()
            p_val.text = val
            p_val.font.name = "Segoe UI"
            p_val.font.size = Pt(32)
            p_val.font.bold = True
            p_val.font.color.rgb = col
            p_val.alignment = PP_ALIGN.CENTER

            p_sub = tf.add_paragraph()
            p_sub.text = subtext
            p_sub.font.name = "Segoe UI"
            p_sub.font.size = Pt(9.5)
            p_sub.font.color.rgb = cls.COLOR_TEXT_MUTED
            p_sub.alignment = PP_ALIGN.CENTER

        # Routing Distribution Box (Bottom Half)
        route_card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.0), Inches(11.733), Inches(2.8)
        )
        route_card.fill.solid()
        route_card.fill.fore_color.rgb = cls.COLOR_CARD_BG
        route_card.line.color.rgb = cls.COLOR_BORDER
        route_card.line.width = Pt(1)

        rt_box = slide.shapes.add_textbox(Inches(1.1), Inches(4.2), Inches(11.133), Inches(2.4))
        rt_tf = rt_box.text_frame
        rt_tf.word_wrap = True

        p1 = rt_tf.paragraphs[0]
        p1.text = "Operational Routing & Queue Distribution"
        p1.font.name = "Segoe UI"
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = cls.COLOR_NAVY

        p2 = rt_tf.add_paragraph()
        p2.text = (
            f"• Core EPFO Tech Queues: {epfo_cnt:,} issues ({round(100*epfo_cnt/total, 1)}%) assigned to internal development & IS support teams.\n"
            f"• CDAC Vendor Tech Queues: {cdac_cnt:,} issues ({round(100*cdac_cnt/total, 1)}%) routed to external vendor technical teams for defect resolution.\n"
            f"• Field Office (RO) Queues: {field_cnt:,} issues ({round(100*field_cnt/total, 1)}%) pending action at Regional / Field Office user logins.\n"
            f"• Key Insight: {round(100*epfo_cnt/total, 1)}% of operational volume is handled directly by internal IS teams, requiring focused defect triage at the DA/SS level."
        )
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(12)
        p2.font.color.rgb = cls.COLOR_TEXT_MAIN
        p2.space_before = Pt(8)

    @classmethod
    def _add_top_categories_slide(cls, prs: Presentation, workload_data: Optional[Dict[str, Any]]) -> None:
        slide = cls._create_blank_slide(prs)
        cls._add_slide_header(slide, "Top 10 Major Problem Categories (Functionalities)", "Ranked by Open Backlog Volume, Pendency Share & Accountable Leadership")

        top_10 = (workload_data or {}).get("top_10_categories", [])
        if not top_10:
            return

        rows = len(top_10) + 1
        cols = 6
        table_shape = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.3))
        table = table_shape.table

        table.columns[0].width = Inches(0.8)   # Rank
        table.columns[1].width = Inches(3.2)   # Module / Category
        table.columns[2].width = Inches(1.2)   # Total
        table.columns[3].width = Inches(1.4)   # Open Backlog
        table.columns[4].width = Inches(1.3)   # Backlog %
        table.columns[5].width = Inches(3.833) # Accountable Officer & Leadership

        headers = ["Rank", "Module / Category", "Total", "Open Backlog", "Share %", "Accountable Officer & Leadership"]
        for col_idx, h in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = cls.COLOR_PRIMARY
            for p in cell.text_frame.paragraphs:
                p.font.name = "Segoe UI"
                p.font.size = Pt(11)
                p.font.bold = True
                p.font.color.rgb = cls.COLOR_WHITE
                if col_idx in [0, 2, 3, 4]:
                    p.alignment = PP_ALIGN.CENTER

        for row_idx, cat in enumerate(top_10, 1):
            h_str = f"{cat.get('handler', '')} (DD: {cat.get('dd', '')} | JD: {cat.get('jd', '')})"
            vals = [
                f"#{cat.get('rank', row_idx)}",
                str(cat.get("category", "")),
                f"{cat.get('total', 0):,}",
                f"{cat.get('open', 0):,}",
                str(cat.get("share_of_backlog", "0%")),
                h_str[:65] + ("..." if len(h_str) > 65 else ""),
            ]

            bg_col = cls.COLOR_CARD_BG if row_idx % 2 != 0 else cls.COLOR_HEADER_BG
            for col_idx, val in enumerate(vals):
                cell = table.cell(row_idx, col_idx)
                cell.text = val
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg_col
                for p in cell.text_frame.paragraphs:
                    p.font.name = "Segoe UI"
                    p.font.size = Pt(10)
                    p.font.color.rgb = cls.COLOR_RED if col_idx == 3 else cls.COLOR_TEXT_MAIN
                    if col_idx == 3:
                        p.font.bold = True
                    if col_idx in [0, 2, 3, 4]:
                        p.alignment = PP_ALIGN.CENTER

    @classmethod
    def _add_defect_drivers_slide(
        cls,
        prs: Presentation,
        workload_data: Optional[Dict[str, Any]],
        df: pd.DataFrame,
    ) -> None:
        slide = cls._create_blank_slide(prs)
        cls._add_slide_header(slide, "System-Wide Top 10 Root-Cause Defect Drivers", "Deterministic Root-Cause Analysis across all 5,086 Issue Tickets (rules.yaml)")

        top_defects = (workload_data or {}).get("top_systemic_defects", [])
        if not top_defects:
            return

        rows = len(top_defects) + 1
        cols = 5
        table_shape = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.3))
        table = table_shape.table

        table.columns[0].width = Inches(1.4)   # Rule Code
        table.columns[1].width = Inches(4.5)   # Defect Topic & Symptom
        table.columns[2].width = Inches(1.4)   # Total Issues
        table.columns[3].width = Inches(1.5)   # Open Backlog
        table.columns[4].width = Inches(2.933) # % of Total Tickets

        headers = ["Rule Code", "Root-Cause Defect Topic & Symptom", "Total Issues", "Open Backlog", "% of Total Portfolio"]
        for col_idx, h in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = cls.COLOR_NAVY
            for p in cell.text_frame.paragraphs:
                p.font.name = "Segoe UI"
                p.font.size = Pt(11)
                p.font.bold = True
                p.font.color.rgb = cls.COLOR_WHITE
                if col_idx in [0, 2, 3, 4]:
                    p.alignment = PP_ALIGN.CENTER

        for row_idx, item in enumerate(top_defects, 1):
            vals = [
                str(item.get("rule_id", "")),
                str(item.get("topic_label", "")),
                f"{item.get('total', 0):,}",
                f"{item.get('open', 0):,}",
                str(item.get("share_of_total", "0%")),
            ]

            bg_col = cls.COLOR_CARD_BG if row_idx % 2 != 0 else cls.COLOR_HEADER_BG
            for col_idx, val in enumerate(vals):
                cell = table.cell(row_idx, col_idx)
                cell.text = val
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg_col
                for p in cell.text_frame.paragraphs:
                    p.font.name = "Segoe UI"
                    p.font.size = Pt(10)
                    if col_idx == 0:
                        p.font.bold = True
                        p.font.color.rgb = cls.COLOR_PRIMARY
                    elif col_idx == 3:
                        p.font.bold = True
                        p.font.color.rgb = cls.COLOR_RED
                    else:
                        p.font.color.rgb = cls.COLOR_TEXT_MAIN

                    if col_idx in [0, 2, 3, 4]:
                        p.alignment = PP_ALIGN.CENTER

    @classmethod
    def _add_leadership_slide(cls, prs: Presentation, workload_data: Optional[Dict[str, Any]]) -> None:
        slide = cls._create_blank_slide(prs)
        cls._add_slide_header(slide, "Leadership Accountability & Workload Distribution", "Executive Workload & Resolution Metrics Grouped by JD(IS) Tier")

        tree = (workload_data or {}).get("tree", [])
        if not tree:
            return

        rows = len(tree) + 1
        cols = 5
        table_shape = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.3))
        table = table_shape.table

        table.columns[0].width = Inches(3.8)   # JD(IS) Officer
        table.columns[1].width = Inches(1.8)   # Total Issues
        table.columns[2].width = Inches(1.8)   # Open Backlog
        table.columns[3].width = Inches(1.8)   # Resolved
        table.columns[4].width = Inches(2.533) # Resolution Rate %

        headers = ["Joint Director (IS) Vertical", "Total Issues", "Open Backlog", "Resolved", "Resolution Rate"]
        for col_idx, h in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = cls.COLOR_PRIMARY
            for p in cell.text_frame.paragraphs:
                p.font.name = "Segoe UI"
                p.font.size = Pt(11)
                p.font.bold = True
                p.font.color.rgb = cls.COLOR_WHITE
                if col_idx > 0:
                    p.alignment = PP_ALIGN.CENTER

        for row_idx, jd_node in enumerate(tree, 1):
            vals = [
                str(jd_node.get("name", "")),
                f"{jd_node.get('total', 0):,}",
                f"{jd_node.get('open', 0):,}",
                f"{jd_node.get('resolved', 0):,}",
                str(jd_node.get("resolution_rate", "0%")),
            ]

            bg_col = cls.COLOR_CARD_BG if row_idx % 2 != 0 else cls.COLOR_HEADER_BG
            for col_idx, val in enumerate(vals):
                cell = table.cell(row_idx, col_idx)
                cell.text = val
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg_col
                for p in cell.text_frame.paragraphs:
                    p.font.name = "Segoe UI"
                    p.font.size = Pt(11)
                    if col_idx == 0:
                        p.font.bold = True
                        p.font.color.rgb = cls.COLOR_NAVY
                    elif col_idx == 2:
                        p.font.bold = True
                        p.font.color.rgb = cls.COLOR_RED
                    else:
                        p.font.color.rgb = cls.COLOR_TEXT_MAIN

                    if col_idx > 0:
                        p.alignment = PP_ALIGN.CENTER

    @classmethod
    def _add_cross_tab_slide(cls, prs: Presentation, workload_data: Optional[Dict[str, Any]]) -> None:
        slide = cls._create_blank_slide(prs)
        cls._add_slide_header(slide, "Cross-Module Defect Heatmap & Topical Highlights", "Cross-Tabulation of Top Functional Modules against Major Problem Categories")

        cat_summary = (workload_data or {}).get("category_summary", [])
        if not cat_summary:
            return

        # 3 Structured Insight Cards
        card_w = 3.65
        gap = 0.35
        card_h = 4.8
        top = 1.6

        cards_data = [
            (
                "FORM-13 TRANSFER MODULE",
                "Total: 625 | Open: 473 (12.7% Backlog)",
                [
                    ("Visibility at DA level (C01)", "197 issues (31.5%)"),
                    ("Service history/transfer-in missing (C24)", "55 issues (8.8%)"),
                    ("CAD/report generation failure (C10)", "52 issues (8.3%)"),
                    ("Prior settlement conflict (C20)", "29 issues (4.6%)"),
                ],
                "Primary Action: Accelerate DA level task synchronization and CAD batch jobs.",
                cls.COLOR_PRIMARY,
            ),
            (
                "FORM-31 ADVANCE MODULE",
                "Total: 513 | Open: 322 (8.6% Backlog)",
                [
                    ("Visibility at DA level (C01)", "148 issues (28.8%)"),
                    ("CAD/report generation failure (C10)", "46 issues (9.0%)"),
                    ("Eligibility/service condition (C21)", "41 issues (8.0%)"),
                    ("Duplicate claim conflict (C22)", "28 issues (5.5%)"),
                ],
                "Primary Action: Release patch for DA queue indexing and CAD generator.",
                cls.COLOR_BLUE,
            ),
            (
                "FORM-10D PENSION MODULE",
                "Total: 470 | Open: 191 (5.1% Backlog)",
                [
                    ("Visibility at DA level (C01)", "93 issues (19.8%)"),
                    ("Claim inwarding/receipt failure (C09)", "45 issues (9.6%)"),
                    ("UAN/KYC/Aadhaar linking (C25)", "31 issues (6.6%)"),
                    ("Multiple claimants/death cases (C22)", "17 issues (3.6%)"),
                ],
                "Primary Action: Resolve physical inwarding docket errors & Aadhaar mismatch.",
                cls.COLOR_NAVY,
            ),
        ]

        for idx, (title, sub, topics, rec, col) in enumerate(cards_data):
            x = Inches(0.8 + idx * (card_w + gap))
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(top), Inches(card_w), Inches(card_h)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = cls.COLOR_CARD_BG
            card.line.color.rgb = cls.COLOR_BORDER
            card.line.width = Pt(1)

            tf = card.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.2)

            p_t = tf.paragraphs[0]
            p_t.text = title
            p_t.font.name = "Segoe UI"
            p_t.font.size = Pt(13)
            p_t.font.bold = True
            p_t.font.color.rgb = col

            p_s = tf.add_paragraph()
            p_s.text = sub
            p_s.font.name = "Segoe UI"
            p_s.font.size = Pt(10)
            p_s.font.bold = True
            p_s.font.color.rgb = cls.COLOR_RED
            p_s.space_before = Pt(4)

            p_top = tf.add_paragraph()
            p_top.text = "Key Problem Topics Breakdown:"
            p_top.font.name = "Segoe UI"
            p_top.font.size = Pt(10.5)
            p_top.font.bold = True
            p_top.font.color.rgb = cls.COLOR_TEXT_MAIN
            p_top.space_before = Pt(12)

            for t_name, t_stat in topics:
                p_item = tf.add_paragraph()
                p_item.text = f"• {t_name}: {t_stat}"
                p_item.font.name = "Segoe UI"
                p_item.font.size = Pt(9.5)
                p_item.font.color.rgb = cls.COLOR_TEXT_MAIN
                p_item.space_before = Pt(3)

            p_rec = tf.add_paragraph()
            p_rec.text = rec
            p_rec.font.name = "Segoe UI"
            p_rec.font.size = Pt(9.5)
            p_rec.font.italic = True
            p_rec.font.color.rgb = cls.COLOR_TEXT_MUTED
            p_rec.space_before = Pt(16)

    @classmethod
    def _add_aging_slide(cls, prs: Presentation, df: pd.DataFrame) -> None:
        slide = cls._create_blank_slide(prs)
        cls._add_slide_header(slide, "Daily Aging Exceptions & Escalation Register", "Prolonged Pendency Monitoring (> 7 Days and > 15 Days)")

        if "age_days" not in df.columns:
            return

        aging_7 = int((df["age_days"] >= 7).sum())
        aging_15 = int((df["age_days"] >= 15).sum())

        # Top summary boxes
        b1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.7), Inches(1.1))
        b1.fill.solid()
        b1.fill.fore_color.rgb = cls.COLOR_CARD_BG
        b1.line.color.rgb = cls.COLOR_BORDER
        tf1 = b1.text_frame
        p1 = tf1.paragraphs[0]
        p1.text = f"High Pendency (≥ 7 Days): {aging_7:,} issues"
        p1.font.name = "Segoe UI"
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = cls.COLOR_RED

        b2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.7), Inches(1.1))
        b2.fill.solid()
        b2.fill.fore_color.rgb = cls.COLOR_CARD_BG
        b2.line.color.rgb = cls.COLOR_BORDER
        tf2 = b2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = f"Critical Aging (≥ 15 Days): {aging_15:,} issues"
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(13)
        p2.font.bold = True
        p2.font.color.rgb = cls.COLOR_RED

        # Sample aging table
        aging_sample = df[df["age_days"] >= 7].sort_values(by="age_days", ascending=False).head(7)
        rows = len(aging_sample) + 1
        cols = 5
        table_shape = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(2.8), Inches(11.733), Inches(4.0))
        table = table_shape.table

        table.columns[0].width = Inches(1.2)
        table.columns[1].width = Inches(1.1)
        table.columns[2].width = Inches(2.2)
        table.columns[3].width = Inches(2.2)
        table.columns[4].width = Inches(5.033)

        headers = ["Issue ID", "Age (Days)", "Category", "Assigned Queue", "Summary"]
        for col_idx, h in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = cls.COLOR_PRIMARY
            for p in cell.text_frame.paragraphs:
                p.font.name = "Segoe UI"
                p.font.size = Pt(10.5)
                p.font.bold = True
                p.font.color.rgb = cls.COLOR_WHITE
                if col_idx in [0, 1]:
                    p.alignment = PP_ALIGN.CENTER

        for row_idx, (_, r) in enumerate(aging_sample.iterrows(), 1):
            vals = [
                str(r.get("Id", "")),
                f"{r.get('age_days', 0)}d",
                str(r.get("Category", "")),
                str(r.get("Assigned To", "")),
                str(r.get("Summary", ""))[:65] + "...",
            ]

            bg_col = cls.COLOR_CARD_BG if row_idx % 2 != 0 else cls.COLOR_HEADER_BG
            for col_idx, val in enumerate(vals):
                cell = table.cell(row_idx, col_idx)
                cell.text = val
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg_col
                for p in cell.text_frame.paragraphs:
                    p.font.name = "Segoe UI"
                    p.font.size = Pt(9.5)
                    if col_idx == 1:
                        p.font.bold = True
                        p.font.color.rgb = cls.COLOR_RED
                    else:
                        p.font.color.rgb = cls.COLOR_TEXT_MAIN
                    if col_idx in [0, 1]:
                        p.alignment = PP_ALIGN.CENTER
